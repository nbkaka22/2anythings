# -*- coding: utf-8 -*-
"""
PDF到PPT转换器 v2.0
使用完全开源的库：pdfplumber + ReportLab + python-pptx
采用图片转换方案，完美保留PDF原始视觉效果
"""

import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
import os
import tempfile
from typing import List, Dict, Tuple, Optional
import io

class PDFToPPTConverterV2:
    """PDF到PPT转换器v2 - 使用完全开源库"""
    
    def __init__(self):
        self.temp_dir = None
        
    def convert_pdf_to_ppt(self, pdf_path: str, output_path: str = None, template_path: str = None) -> str:
        """
        将PDF转换为PPT（使用图片转换方案）
        
        Args:
            pdf_path: PDF文件路径
            output_path: 输出PPT文件路径，如果为None则自动生成
            template_path: PPT模板文件路径，如果为None则使用默认模板
            
        Returns:
            输出文件路径
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
            
        if output_path is None:
            base_name = os.path.splitext(os.path.basename(pdf_path))[0]
            output_path = os.path.join(os.path.dirname(pdf_path), f"{base_name}_converted_v2.pptx")
            
        # 创建临时目录
        self.temp_dir = tempfile.mkdtemp()
        
        try:
            # 创建空白PPT演示文稿
            ppt = Presentation()
            print("使用默认空白演示文稿")
            
            # 使用pdfplumber打开PDF
            with pdfplumber.open(pdf_path) as pdf:
                total_pages = len(pdf.pages)
                print(f"开始转换PDF，共{total_pages}页...")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    print(f"正在处理第{page_num}/{total_pages}页...")
                    
                    # 将PDF页面转换为图片
                    page_image = self._convert_page_to_image(page, page_num)
                    
                    # 添加新幻灯片
                    slide_layout = ppt.slide_layouts[6]  # 空白布局
                    slide = ppt.slides.add_slide(slide_layout)
                    
                    # 将图片添加到幻灯片
                    self._add_image_to_slide(slide, page_image, ppt)
                    
            # 保存PPT文件
            ppt.save(output_path)
            print(f"转换完成！输出文件：{output_path}")
            
            return output_path
            
        finally:
            # 清理临时文件
            self._cleanup_temp_files()
    
    def _convert_page_to_image(self, page, page_num: int) -> str:
        """
        将PDF页面转换为高质量图片
        
        Args:
            page: pdfplumber页面对象
            page_num: 页面编号
            
        Returns:
            图片文件路径
        """
        try:
            # 使用pdfplumber的内置方法转换为图片
            # 设置较高的分辨率以保证质量
            img = page.to_image(resolution=200)
            
            # 保存为临时文件
            image_path = os.path.join(self.temp_dir, f"page_{page_num}.png")
            img.save(image_path, format='PNG')
            
            return image_path
            
        except Exception as e:
            print(f"页面{page_num}转换为图片时出错: {e}")
            # 创建一个错误提示图片
            return self._create_error_image(page_num)
    
    def _create_error_image(self, page_num: int) -> str:
        """
        创建错误提示图片
        
        Args:
            page_num: 页面编号
            
        Returns:
            错误图片文件路径
        """
        # 创建一个简单的错误提示图片
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)
        
        error_text = f"页面 {page_num} 转换失败"
        # 使用默认字体
        draw.text((400, 300), error_text, fill='red', anchor='mm')
        
        # 保存错误图片
        error_image_path = os.path.join(self.temp_dir, f"error_page_{page_num}.png")
        img.save(error_image_path)
        
        return error_image_path
    
    def _add_image_to_slide(self, slide, image_path: str, presentation):
        """
        将图片添加到幻灯片
        
        Args:
            slide: PPT幻灯片对象
            image_path: 图片文件路径
            presentation: PPT演示文稿对象
        """
        try:
            # 获取幻灯片尺寸（使用presentation对象）
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            
            # 打开图片获取原始尺寸
            with Image.open(image_path) as img:
                img_width, img_height = img.size
            
            # 计算缩放比例，保持宽高比
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale_ratio = min(width_ratio, height_ratio)
            
            # 计算最终尺寸
            final_width = int(img_width * scale_ratio)
            final_height = int(img_height * scale_ratio)
            
            # 计算居中位置
            left = (slide_width - final_width) // 2
            top = (slide_height - final_height) // 2
            
            # 添加图片到幻灯片
            slide.shapes.add_picture(
                image_path,
                left, top,
                final_width, final_height
            )
            
        except Exception as e:
            print(f"添加图片到幻灯片时出错: {e}")
    
    def _cleanup_temp_files(self):
        """
        清理临时文件
        """
        if self.temp_dir and os.path.exists(self.temp_dir):
            try:
                import shutil
                shutil.rmtree(self.temp_dir)
            except Exception as e:
                print(f"清理临时文件时出错: {e}")
    
    def get_pdf_info(self, pdf_path: str) -> Dict:
        """
        获取PDF文件信息
        
        Args:
            pdf_path: PDF文件路径
            
        Returns:
            PDF信息字典
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
        
        try:
            with pdfplumber.open(pdf_path) as pdf:
                info = {
                    'pages': len(pdf.pages),
                    'metadata': pdf.metadata or {},
                    'first_page_size': None
                }
                
                if pdf.pages:
                    first_page = pdf.pages[0]
                    info['first_page_size'] = {
                        'width': first_page.width,
                        'height': first_page.height
                    }
                
                return info
                
        except Exception as e:
            return {
                'error': f"读取PDF信息时出错: {e}",
                'pages': 0,
                'metadata': {},
                'first_page_size': None
            }

# 兼容性函数，保持与原版本的接口一致
def convert_pdf_to_ppt(pdf_path: str, output_path: str = None, template_path: str = None) -> str:
    """
    PDF到PPT转换的便捷函数
    
    Args:
        pdf_path: PDF文件路径
        output_path: 输出PPT文件路径
        template_path: PPT模板文件路径
        
    Returns:
        输出文件路径
    """
    converter = PDFToPPTConverterV2()
    return converter.convert_pdf_to_ppt(pdf_path, output_path, template_path)

if __name__ == "__main__":
    # 测试代码
    import sys
    
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
        
        try:
            result = convert_pdf_to_ppt(pdf_file, output_file)
            print(f"转换成功！输出文件：{result}")
        except Exception as e:
            print(f"转换失败：{e}")
    else:
        print("使用方法：python pdf_to_ppt_converter_v2.py <PDF文件路径> [输出文件路径]")