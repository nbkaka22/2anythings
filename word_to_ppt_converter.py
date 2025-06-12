"""Word到PPT转换器
通过读取Word文档内容，智能转换为PPT格式
"""

import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import tempfile
import logging

logger = logging.getLogger('pdf_converter')

class WordToPPTConverter:
    """Word文档到PPT转换器"""
    
    def __init__(self):
        self.max_chars_per_slide = 800  # 每页最大字符数
        self.max_lines_per_slide = 15   # 每页最大行数
        
        # 预定义样式主题
        self.theme_colors = {
            'primary': RGBColor(68, 114, 196),      # 蓝色主题
            'secondary': RGBColor(112, 173, 71),    # 绿色辅助
            'accent': RGBColor(255, 192, 0),        # 橙色强调
            'text': RGBColor(68, 68, 68),           # 深灰文本
            'background': RGBColor(248, 249, 250)   # 浅灰背景
        }
        
        # 字体设置
        self.fonts = {
            'title': '微软雅黑',
            'heading': '微软雅黑',
            'body': '微软雅黑'
        }
    
    def convert_word_to_ppt(self, word_path, output_path, template_path=None):
        """将Word文档转换为PPT
        
        Args:
            word_path: Word文档路径
            output_path: 输出PPT路径
            template_path: PPT模板文件路径，如果为None则使用默认模板
            
        Returns:
            输出文件路径
        """
        try:
            # 读取Word文档
            doc = Document(word_path)
            
            # 创建空白演示文稿
            ppt = Presentation()
            logger.info("使用默认空白样式创建PPT")
            
            # 提取Word内容
            content_blocks = self._extract_content_from_word(doc)
            
            # 将内容转换为PPT页面
            self._create_ppt_slides(ppt, content_blocks)
            
            # 保存PPT
            ppt.save(output_path)
            
            logger.info(f"Word到PPT转换成功: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Word到PPT转换失败: {str(e)}")
            raise e
    
    def _extract_content_from_word(self, doc):
        """从Word文档中提取内容
        
        Returns:
            list: 内容块列表，每个块包含类型和内容
        """
        content_blocks = []
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
                
            # 判断段落类型
            if self._is_heading(paragraph):
                content_blocks.append({
                    'type': 'heading',
                    'text': text,
                    'level': self._get_heading_level(paragraph)
                })
            else:
                content_blocks.append({
                    'type': 'paragraph',
                    'text': text
                })
        
        return content_blocks
    
    def _is_heading(self, paragraph):
        """判断段落是否为标题"""
        style_name = paragraph.style.name.lower()
        return 'heading' in style_name or paragraph.runs and paragraph.runs[0].bold
    
    def _get_heading_level(self, paragraph):
        """获取标题级别"""
        style_name = paragraph.style.name.lower()
        if 'heading 1' in style_name:
            return 1
        elif 'heading 2' in style_name:
            return 2
        elif 'heading 3' in style_name:
            return 3
        else:
            return 1  # 默认为一级标题
    
    def _create_ppt_slides(self, ppt, content_blocks):
        """根据内容块创建PPT幻灯片"""
        if not content_blocks:
            return
        
        current_slide = None
        current_content = []
        current_chars = 0
        
        for block in content_blocks:
            # 如果是一级标题，创建新幻灯片
            if block['type'] == 'heading' and block['level'] == 1:
                # 保存当前幻灯片内容
                if current_slide is not None and current_content:
                    self._add_content_to_slide(current_slide, current_content)
                
                # 创建新幻灯片
                slide_layout = ppt.slide_layouts[1]  # 使用标题和内容布局
                current_slide = ppt.slides.add_slide(slide_layout)
                
                # 设置标题
                title = current_slide.shapes.title
                # 使用样式系统设置标题
                self._apply_text_style(title.text_frame, block['text'], 'title')
                
                # 重置内容
                current_content = []
                current_chars = 0
                
            else:
                # 检查是否需要分页
                block_chars = len(block['text'])
                if (current_chars + block_chars > self.max_chars_per_slide or 
                    len(current_content) >= self.max_lines_per_slide):
                    
                    # 保存当前幻灯片内容
                    if current_slide is not None and current_content:
                        self._add_content_to_slide(current_slide, current_content)
                    
                    # 创建新幻灯片
                    slide_layout = ppt.slide_layouts[1]
                    current_slide = ppt.slides.add_slide(slide_layout)
                    
                    # 设置默认标题
                    title = current_slide.shapes.title
                    self._apply_text_style(title.text_frame, "内容续页", 'heading')
                    
                    # 重置内容
                    current_content = []
                    current_chars = 0
                
                # 添加内容到当前页
                current_content.append(block)
                current_chars += block_chars
        
        # 处理最后一页
        if current_slide is not None and current_content:
            self._add_content_to_slide(current_slide, current_content)
        
        # 如果没有创建任何幻灯片，创建一个默认幻灯片
        if len(ppt.slides) == 0:
            slide_layout = ppt.slide_layouts[0]  # 标题幻灯片
            slide = ppt.slides.add_slide(slide_layout)
            title = slide.shapes.title
            self._apply_text_style(title.text_frame, "PDF转换内容", 'title')
            subtitle = slide.placeholders[1]
            self._apply_text_style(subtitle.text_frame, "通过Word中转生成的PPT", 'body')
    
    def _add_content_to_slide(self, slide, content_blocks):
        """将内容添加到幻灯片"""
        if not content_blocks:
            return
        
        # 优先查找内容占位符
        content_placeholder = None
        
        # 尝试多种方式查找内容占位符
        for placeholder in slide.placeholders:
            # 方法1：通过占位符类型查找（内容占位符）
            if placeholder.placeholder_format.type == 2:
                content_placeholder = placeholder
                break
            # 方法2：通过索引查找（通常索引1是内容占位符）
            elif placeholder.placeholder_format.idx == 1:
                content_placeholder = placeholder
                break
        
        # 如果还没找到，尝试查找任何可用的文本占位符
        if content_placeholder is None:
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'text_frame') and placeholder != slide.shapes.title:
                    content_placeholder = placeholder
                    break
        
        # 最后才考虑创建新文本框
        if content_placeholder is None:
            # 如果没有找到任何内容占位符，创建文本框
            logger.info("未找到内容占位符，创建新文本框")
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
        else:
            logger.info(f"使用默认内容占位符，类型: {content_placeholder.placeholder_format.type}, 索引: {content_placeholder.placeholder_format.idx}")
            text_frame = content_placeholder.text_frame
        
        # 优化文本布局
        self._optimize_text_layout(text_frame)
        
        # 清空现有内容
        text_frame.clear()
        
        # 添加内容
        for i, block in enumerate(content_blocks):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            run = p.add_run()
            run.text = block['text']
            
            # 使用新的样式系统设置格式
            if block['type'] == 'heading':
                run.font.name = self.fonts['heading']
                run.font.bold = True
                run.font.color.rgb = self.theme_colors['secondary']
                if block['level'] == 2:
                    run.font.size = Pt(18)
                elif block['level'] == 3:
                    run.font.size = Pt(16)
                else:
                    run.font.size = Pt(20)
                p.alignment = PP_ALIGN.LEFT
            else:
                run.font.name = self.fonts['body']
                run.font.size = Pt(14)
                run.font.color.rgb = self.theme_colors['text']
                p.alignment = PP_ALIGN.LEFT
            
            # 设置行间距和段落间距
            p.space_after = Pt(8)
            p.space_before = Pt(2)
    
    def _apply_slide_theme(self, slide):
        """为幻灯片应用主题样式"""
        try:
            # 设置幻灯片背景（如果支持）
            # 注意：python-pptx对背景设置的支持有限
            pass
        except Exception as e:
            logger.warning(f"应用幻灯片主题时出错: {e}")
    
    def _optimize_text_layout(self, text_frame):
        """优化文本布局"""
        try:
            # 设置文本框属性
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            text_frame.word_wrap = True
            
            # 自动调整文本大小
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
        except Exception as e:
            logger.warning(f"优化文本布局时出错: {e}")
    
    def _apply_text_style(self, text_frame, text_content, style_type='body'):
        """应用文本样式"""
        try:
            # 优化文本布局
            self._optimize_text_layout(text_frame)
            
            # 清空现有内容
            text_frame.clear()
            
            # 添加段落
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text_content
            
            # 根据样式类型设置格式
            if style_type == 'title':
                run.font.name = self.fonts['title']
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = self.theme_colors['primary']
                p.alignment = PP_ALIGN.CENTER
            elif style_type == 'heading':
                run.font.name = self.fonts['heading']
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = self.theme_colors['secondary']
                p.alignment = PP_ALIGN.LEFT
            else:  # body
                run.font.name = self.fonts['body']
                run.font.size = Pt(14)
                run.font.color.rgb = self.theme_colors['text']
                p.alignment = PP_ALIGN.LEFT
                
            # 设置行间距
            p.space_after = Pt(6)
            
        except Exception as e:
            logger.warning(f"应用文本样式时出错: {e}")
            # 回退到简单文本设置
            text_frame.text = text_content
    
    def _is_title(self, text):
        """判断是否为标题"""
        text = text.strip()
        if len(text) == 0:
            return False
            
        # 标题特征：较短、包含关键词、或者格式特殊
        title_keywords = ['第', '章', '节', '部分', '概述', '介绍', '总结', '结论']
        
        # 长度判断
        if len(text) <= 50 and any(keyword in text for keyword in title_keywords):
            return True
            
        # 数字开头的标题
        if len(text) <= 100 and (text[0].isdigit() or text.startswith(('一', '二', '三', '四', '五'))):
            return True
            
        return False