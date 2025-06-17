"""Word转PPT转换器插件

将Word文档转换为PowerPoint演示文稿格式
"""

import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging
from typing import Dict, Any, List
import re

# 导入基类
from converters.converter_interface import ConverterInterface, ConverterMetadata

logger = logging.getLogger('pdf_converter')

class WordToPPTConverter(ConverterInterface):
    """Word转PPT转换器
    
    将Word文档的内容转换为PowerPoint演示文稿
    """
    
    def __init__(self):
        self._temp_files = []
    
    @property
    def name(self) -> str:
        return "word_to_ppt"
    
    @property
    def description(self) -> str:
        return "将Word文档转换为PowerPoint演示文稿，智能识别标题和内容"
    
    @property
    def supported_input_formats(self) -> List[str]:
        return ["docx", "doc"]
    
    @property
    def supported_output_formats(self) -> List[str]:
        return ["pptx"]
    
    @property
    def version(self) -> str:
        return "1.0.0"
    
    def validate_input(self, input_path: str) -> bool:
        """验证Word文件是否有效"""
        try:
            if not os.path.exists(input_path):
                logger.error(f"文件不存在: {input_path}")
                return False
            
            if not input_path.lower().endswith(('.docx', '.doc')):
                logger.error(f"文件格式不正确，需要Word文件: {input_path}")
                return False
            
            # 尝试打开Word文档
            if input_path.lower().endswith('.docx'):
                doc = Document(input_path)
                if len(doc.paragraphs) == 0:
                    logger.error(f"Word文档为空: {input_path}")
                    return False
            else:
                # .doc文件需要特殊处理
                logger.warning(f"检测到.doc文件，建议转换为.docx格式: {input_path}")
            
            return True
            
        except Exception as e:
            logger.error(f"Word文件验证失败: {e}")
            return False
    
    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        """执行Word到PPT的转换
        
        Args:
            input_path: Word文件路径
            output_path: 输出PPT文件路径
            **kwargs: 转换选项
                - slide_layout: 幻灯片布局类型 ('title_content', 'title_only', 'content_only')
                - max_content_per_slide: 每张幻灯片最大内容行数
                - font_size_title: 标题字体大小
                - font_size_content: 内容字体大小
                - auto_split: 是否自动分割长内容
        
        Returns:
            bool: 转换是否成功
        """
        try:
            # 获取转换选项
            slide_layout = kwargs.get('slide_layout', 'title_content')
            max_content_per_slide = kwargs.get('max_content_per_slide', 10)
            font_size_title = kwargs.get('font_size_title', 32)
            font_size_content = kwargs.get('font_size_content', 18)
            auto_split = kwargs.get('auto_split', True)
            
            logger.info(f"开始Word转PPT转换: {input_path} -> {output_path}")
            logger.info(f"布局: {slide_layout}, 最大内容行数: {max_content_per_slide}")
            
            # 确保输出目录存在
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 只支持.docx文件
            if not input_path.lower().endswith('.docx'):
                logger.error("当前只支持.docx格式的Word文档")
                return False
            
            # 读取Word文档
            doc = Document(input_path)
            
            # 解析文档结构
            slides_data = self._parse_document_structure(doc, max_content_per_slide, auto_split)
            
            if not slides_data:
                logger.error("未能从Word文档中提取有效内容")
                return False
            
            # 创建PowerPoint演示文稿
            prs = Presentation()
            
            # 生成幻灯片
            for slide_data in slides_data:
                self._create_slide(prs, slide_data, slide_layout, font_size_title, font_size_content)
            
            # 保存PowerPoint文件
            prs.save(output_path)
            
            # 验证输出文件
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"Word转PPT成功: {output_path}, 共生成 {len(slides_data)} 张幻灯片")
                return True
            else:
                logger.error("转换完成但输出文件无效")
                return False
                
        except Exception as e:
            logger.error(f"Word转PPT失败: {e}")
            return False
    
    def _parse_document_structure(self, doc: Document, max_content_per_slide: int, auto_split: bool) -> List[Dict[str, Any]]:
        """解析Word文档结构"""
        slides_data = []
        current_slide = {'title': '', 'content': []}
        
        try:
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # 判断是否为标题
                if self._is_title(paragraph, text):
                    # 如果当前幻灯片有内容，保存它
                    if current_slide['title'] or current_slide['content']:
                        slides_data.append(current_slide)
                    
                    # 开始新幻灯片
                    current_slide = {'title': text, 'content': []}
                else:
                    # 添加到内容
                    current_slide['content'].append(text)
                    
                    # 如果启用自动分割且内容过多
                    if auto_split and len(current_slide['content']) >= max_content_per_slide:
                        slides_data.append(current_slide)
                        # 继续使用相同标题创建新幻灯片
                        current_slide = {
                            'title': current_slide['title'] + ' (续)',
                            'content': []
                        }
            
            # 添加最后一张幻灯片
            if current_slide['title'] or current_slide['content']:
                slides_data.append(current_slide)
            
            # 如果没有找到标题，创建默认结构
            if not slides_data:
                all_content = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                if all_content:
                    slides_data.append({
                        'title': '文档内容',
                        'content': all_content[:max_content_per_slide]
                    })
            
            logger.info(f"解析完成，共识别 {len(slides_data)} 张幻灯片")
            return slides_data
            
        except Exception as e:
            logger.error(f"文档结构解析失败: {e}")
            return []
    
    def _is_title(self, paragraph, text: str) -> bool:
        """判断段落是否为标题"""
        try:
            # 检查样式名称
            style_name = paragraph.style.name.lower()
            if 'heading' in style_name or 'title' in style_name:
                return True
            
            # 检查字体大小（标题通常字体较大）
            if paragraph.runs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt > 14:
                        return True
            
            # 检查文本特征
            # 短文本且不以句号结尾
            if len(text) < 100 and not text.endswith(('。', '.', '!', '?')):
                # 检查是否包含数字编号
                if re.match(r'^\d+[.、]', text) or re.match(r'^[一二三四五六七八九十]+[、.]', text):
                    return True
                
                # 检查是否全大写或包含特殊格式
                if text.isupper() or text.startswith(('第', '章', '节')):
                    return True
            
            return False
            
        except Exception as e:
            logger.warning(f"标题判断失败: {e}")
            return False
    
    def _create_slide(self, prs: Presentation, slide_data: Dict[str, Any], 
                     layout_type: str, font_size_title: int, font_size_content: int):
        """创建幻灯片"""
        try:
            # 选择布局
            if layout_type == 'title_only':
                layout = prs.slide_layouts[5]  # 标题幻灯片
            elif layout_type == 'content_only':
                layout = prs.slide_layouts[6]  # 空白幻灯片
            else:  # title_content
                layout = prs.slide_layouts[1]  # 标题和内容
            
            slide = prs.slides.add_slide(layout)
            
            # 添加标题
            if slide_data['title'] and hasattr(slide.shapes, 'title'):
                title_shape = slide.shapes.title
                title_shape.text = slide_data['title']
                
                # 设置标题格式
                title_frame = title_shape.text_frame
                title_frame.paragraphs[0].font.size = Pt(font_size_title)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 添加内容
            if slide_data['content']:
                if layout_type == 'title_content' and len(slide.placeholders) > 1:
                    # 使用内容占位符
                    content_shape = slide.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.clear()
                    
                    for i, content_line in enumerate(slide_data['content']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = content_line
                        p.font.size = Pt(font_size_content)
                        p.level = 0
                else:
                    # 手动添加文本框
                    left = Inches(1)
                    top = Inches(2) if slide_data['title'] else Inches(1)
                    width = Inches(8)
                    height = Inches(5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    
                    for i, content_line in enumerate(slide_data['content']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = content_line
                        p.font.size = Pt(font_size_content)
            
        except Exception as e:
            logger.error(f"创建幻灯片失败: {e}")
    
    def get_default_options(self) -> Dict[str, Any]:
        """获取默认转换选项"""
        return {
            'slide_layout': 'title_content',
            'max_content_per_slide': 10,
            'font_size_title': 32,
            'font_size_content': 18,
            'auto_split': True,
            'preserve_formatting': False
        }
    
    def cleanup(self):
        """清理临时文件"""
        for temp_file in self._temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.warning(f"清理临时文件失败: {temp_file}, {e}")
        
        self._temp_files.clear()
    
    def _extract_tables(self, doc: Document) -> List[List[List[str]]]:
        """提取Word文档中的表格"""
        tables_data = []
        
        try:
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                tables_data.append(table_data)
            
            logger.info(f"提取到 {len(tables_data)} 个表格")
            return tables_data
            
        except Exception as e:
            logger.warning(f"表格提取失败: {e}")
            return []
    
    def _add_table_to_slide(self, slide, table_data: List[List[str]]):
        """将表格添加到幻灯片"""
        try:
            if not table_data:
                return
            
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            
            if rows == 0 or cols == 0:
                return
            
            # 添加表格
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(4)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # 填充表格数据
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    if j < len(table.columns):
                        cell = table.cell(i, j)
                        cell.text = cell_data
                        
                        # 设置字体
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
            
        except Exception as e:
            logger.warning(f"添加表格失败: {e}")

# 转换器元数据
WORD_TO_PPT_METADATA = ConverterMetadata(
    name="word_to_ppt",
    description="将Word文档转换为PowerPoint演示文稿，智能识别标题和内容结构",
    version="1.0.0",
    author="PDF Converter Team",
    supported_input_formats=["docx", "doc"],
    supported_output_formats=["pptx"],
    dependencies=["python-docx", "python-pptx"],
    priority=10
)