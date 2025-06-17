"""PDF转PPT转换器插件

将PDF文件转换为PowerPoint演示文稿格式
"""

import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
import logging
from typing import Dict, Any, List
import io

# 导入基类
from converters.converter_interface import ConverterInterface, ConverterMetadata

logger = logging.getLogger('pdf_converter')

class PDFToPPTConverter(ConverterInterface):
    """PDF转PPT转换器
    
    将PDF文件的每一页转换为PPT幻灯片
    """
    
    def __init__(self):
        self._temp_files = []
    
    @property
    def name(self) -> str:
        return "pdf_to_ppt"
    
    @property
    def description(self) -> str:
        return "将PDF文件转换为PowerPoint演示文稿，每页PDF对应一张幻灯片"
    
    @property
    def supported_input_formats(self) -> List[str]:
        return ["pdf"]
    
    @property
    def supported_output_formats(self) -> List[str]:
        return ["pptx"]
    
    @property
    def version(self) -> str:
        return "1.0.0"
    
    def validate_input(self, input_path: str) -> bool:
        """验证PDF文件是否有效"""
        try:
            if not os.path.exists(input_path):
                logger.error(f"文件不存在: {input_path}")
                return False
            
            if not input_path.lower().endswith('.pdf'):
                logger.error(f"文件格式不正确，需要PDF文件: {input_path}")
                return False
            
            # 尝试打开PDF文件
            doc = fitz.open(input_path)
            if doc.page_count == 0:
                logger.error(f"PDF文件为空: {input_path}")
                doc.close()
                return False
            
            doc.close()
            return True
            
        except Exception as e:
            logger.error(f"PDF文件验证失败: {e}")
            return False
    
    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        """执行PDF到PPT的转换
        
        Args:
            input_path: PDF文件路径
            output_path: 输出PPT文件路径
            **kwargs: 转换选项
                - dpi: 图像分辨率 (默认150)
                - image_format: 图像格式 ('png' 或 'jpeg')
                - start_page: 起始页码（从0开始）
                - end_page: 结束页码
                - slide_layout: 幻灯片布局索引
                - include_text: 是否包含文本内容
        
        Returns:
            bool: 转换是否成功
        """
        try:
            # 获取转换选项
            dpi = kwargs.get('dpi', 150)
            image_format = kwargs.get('image_format', 'png')
            start_page = kwargs.get('start_page', 0)
            end_page = kwargs.get('end_page', None)
            slide_layout = kwargs.get('slide_layout', 6)  # 空白布局
            include_text = kwargs.get('include_text', True)
            
            logger.info(f"开始PDF转PPT转换: {input_path} -> {output_path}")
            logger.info(f"DPI: {dpi}, 格式: {image_format}, 包含文本: {include_text}")
            
            # 确保输出目录存在
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 打开PDF文档
            pdf_doc = fitz.open(input_path)
            
            if end_page is None:
                end_page = pdf_doc.page_count - 1
            
            # 创建PowerPoint演示文稿
            prs = Presentation()
            
            # 逐页转换
            for page_num in range(start_page, min(end_page + 1, pdf_doc.page_count)):
                logger.info(f"转换第 {page_num + 1} 页")
                
                page = pdf_doc[page_num]
                
                # 添加新幻灯片
                slide_layout_obj = prs.slide_layouts[slide_layout]
                slide = prs.slides.add_slide(slide_layout_obj)
                
                # 将PDF页面转换为图像
                if self._add_page_as_image(page, slide, dpi, image_format):
                    # 如果需要，添加文本内容
                    if include_text:
                        self._add_text_content(page, slide)
                else:
                    logger.warning(f"第 {page_num + 1} 页图像转换失败")
            
            pdf_doc.close()
            
            # 保存PowerPoint文件
            prs.save(output_path)
            
            # 验证输出文件
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"PDF转PPT成功: {output_path}")
                return True
            else:
                logger.error("转换完成但输出文件无效")
                return False
                
        except Exception as e:
            logger.error(f"PDF转PPT失败: {e}")
            return False
    
    def _add_page_as_image(self, page, slide, dpi: int, image_format: str) -> bool:
        """将PDF页面作为图像添加到幻灯片"""
        try:
            # 将PDF页面渲染为图像
            mat = fitz.Matrix(dpi/72, dpi/72)  # 缩放矩阵
            pix = page.get_pixmap(matrix=mat)
            
            # 转换为PIL图像
            img_data = pix.tobytes(image_format)
            img = Image.open(io.BytesIO(img_data))
            
            # 创建临时图像文件
            temp_img_path = tempfile.mktemp(suffix=f'.{image_format}')
            self._temp_files.append(temp_img_path)
            img.save(temp_img_path)
            
            # 计算图像在幻灯片中的位置和大小
            slide_width = Inches(10)  # 标准幻灯片宽度
            slide_height = Inches(7.5)  # 标准幻灯片高度
            
            # 计算缩放比例以适应幻灯片
            img_width, img_height = img.size
            width_ratio = slide_width.inches / (img_width / dpi)
            height_ratio = slide_height.inches / (img_height / dpi)
            scale_ratio = min(width_ratio, height_ratio, 1.0)  # 不放大
            
            # 计算最终尺寸和位置
            final_width = Inches((img_width / dpi) * scale_ratio)
            final_height = Inches((img_height / dpi) * scale_ratio)
            left = (slide_width - final_width) / 2
            top = (slide_height - final_height) / 2
            
            # 添加图像到幻灯片
            slide.shapes.add_picture(temp_img_path, left, top, final_width, final_height)
            
            return True
            
        except Exception as e:
            logger.error(f"添加页面图像失败: {e}")
            return False
    
    def _add_text_content(self, page, slide):
        """添加文本内容到幻灯片"""
        try:
            # 提取页面文本
            text = page.get_text().strip()
            
            if text:
                # 在幻灯片底部添加文本框
                left = Inches(0.5)
                top = Inches(6.5)
                width = Inches(9)
                height = Inches(1)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text[:200] + "..." if len(text) > 200 else text  # 限制文本长度
                
                # 设置文本格式
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Inches(0.1)  # 小字体
                    paragraph.font.name = 'Arial'
                
        except Exception as e:
            logger.warning(f"添加文本内容失败: {e}")
    
    def get_default_options(self) -> Dict[str, Any]:
        """获取默认转换选项"""
        return {
            'dpi': 150,
            'image_format': 'png',
            'start_page': 0,
            'end_page': None,
            'slide_layout': 6,  # 空白布局
            'include_text': False,  # 默认不包含文本
            'preserve_aspect_ratio': True
        }
    
    def cleanup(self):
        """清理临时文件"""
        for temp_file in self._temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.warning(f"清理临时文件失败: {temp_file}, {e}")
        
        self._temp_files.clear()
    
    def _optimize_image_quality(self, img: Image.Image, target_size_mb: float = 1.0) -> Image.Image:
        """优化图像质量和大小"""
        try:
            # 计算当前图像大小（估算）
            width, height = img.size
            current_size_mb = (width * height * 3) / (1024 * 1024)  # RGB估算
            
            if current_size_mb > target_size_mb:
                # 计算缩放比例
                scale_factor = (target_size_mb / current_size_mb) ** 0.5
                new_width = int(width * scale_factor)
                new_height = int(height * scale_factor)
                
                # 使用高质量重采样
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                logger.info(f"图像已优化: {width}x{height} -> {new_width}x{new_height}")
            
            return img
            
        except Exception as e:
            logger.warning(f"图像优化失败: {e}")
            return img

# 转换器元数据
PDF_TO_PPT_METADATA = ConverterMetadata(
    name="pdf_to_ppt",
    description="将PDF文件转换为PowerPoint演示文稿，每页PDF对应一张幻灯片",
    version="1.0.0",
    author="PDF Converter Team",
    supported_input_formats=["pdf"],
    supported_output_formats=["pptx"],
    dependencies=["PyMuPDF", "python-pptx", "Pillow"],
    priority=10
)